"""Text extraction from uploaded documents.

Dispatches to format-specific extractors based on file extension.
"""

from __future__ import annotations

from collections.abc import Callable
from pathlib import Path

import pdfplumber
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation as PptxPresentation
from striprtf.striprtf import rtf_to_text

from ananta.explorers.document.config import MAX_EXTRACTED_TEXT_BYTES

# Extensions treated as plain text (read with open()).
#
# `.env` is intentionally NOT in this list (I8): files named `.env` typically
# contain API keys, DB credentials, and other secrets. A user who drops a
# project folder onto the explorer should not silently land their secrets
# in the RLM document store, where the LLM reads them and the LLM provider
# sees them. Combined with [C2], this also closes the drive-by-credential-
# exfiltration vector. If a user genuinely wants to ingest a config file,
# they can rename it to `.env.txt` or use a different non-secret extension.
_PLAIN_TEXT_EXTENSIONS = frozenset(
    {
        ".txt",
        ".md",
        ".csv",
        ".log",
        ".json",
        ".yaml",
        ".yml",
        ".xml",
        ".html",
        ".htm",
        ".ini",
        ".cfg",
        ".toml",
        ".py",
        ".js",
        ".ts",
        ".java",
        ".c",
        ".cpp",
        ".h",
        ".rs",
        ".go",
        ".rb",
        ".sh",
        ".bat",
        ".sql",
        ".r",
        ".tex",
    }
)


_SUPPORTED_EXTENSIONS = _PLAIN_TEXT_EXTENSIONS | frozenset(
    {".pdf", ".docx", ".pptx", ".xlsx", ".rtf"}
)


def is_supported_extension(filename: str) -> bool:
    """Check if a filename has a supported extension for text extraction."""
    ext = Path(filename).suffix.lower()
    return ext in _SUPPORTED_EXTENSIONS


def _truncate_to_cap(text: str) -> str:
    """Truncate *text* to ``MAX_EXTRACTED_TEXT_BYTES`` UTF-8 bytes (I3).

    Bounds in-process memory and the on-disk document store against
    decompression-bomb DoS via heavily-compressed xlsx/pptx/docx/pdf,
    and reduces attacker-controlled bytes flowing into the RLM context.
    """
    # Fast path: each UTF-8 char is at most 4 bytes, so a char count well
    # under the cap can't possibly encode to more bytes than the cap. Skips
    # an unnecessary bytes-copy for the common (under-cap) case.
    if len(text) * 4 <= MAX_EXTRACTED_TEXT_BYTES:
        return text
    encoded = text.encode("utf-8")
    if len(encoded) <= MAX_EXTRACTED_TEXT_BYTES:
        return text
    # errors='ignore' drops a partial multi-byte sequence at the boundary
    # so the result is always valid UTF-8 with no replacement characters.
    truncated = encoded[:MAX_EXTRACTED_TEXT_BYTES].decode("utf-8", errors="ignore")
    return (
        truncated + f"\n[Extracted text truncated to {MAX_EXTRACTED_TEXT_BYTES:,} of "
        f"{len(encoded):,} bytes.]"
    )


def extract_text(path: Path, content_type: str | None = None) -> str:
    """Extract text content from a file.

    *content_type* is accepted for forward-compatibility but dispatch
    is by file extension. Raises ``ValueError`` for unsupported types
    or for corrupt/unreadable files (translated from per-library errors
    like pdfplumber's PDFSyntaxError, zipfile's BadZipFile, openpyxl's
    InvalidFileException, python-pptx's PackageNotFoundError).

    Output is capped at ``MAX_EXTRACTED_TEXT_BYTES`` UTF-8 bytes with an
    inline truncation marker — see ``_truncate_to_cap``.
    """
    ext = path.suffix.lower()

    if ext in _PLAIN_TEXT_EXTENSIONS:
        return _truncate_to_cap(_extract_plain_text(path))
    fmt_extractors: dict[str, tuple[str, Callable[[Path], str]]] = {
        ".pdf": ("pdf", _extract_pdf),
        ".docx": ("docx", _extract_docx),
        ".pptx": ("pptx", _extract_pptx),
        ".xlsx": ("xlsx", _extract_xlsx),
        ".rtf": ("rtf", _extract_rtf),
    }
    if ext in fmt_extractors:
        fmt_name, fn = fmt_extractors[ext]
        try:
            return _truncate_to_cap(fn(path))
        except ValueError:
            raise
        except Exception as exc:
            # Per-format libraries raise their own exception types on corrupt
            # input (e.g. zipfile.BadZipFile). Translate to ValueError so the
            # API's per-file ValueError handler emits a "text extraction
            # failed" reason instead of a generic "unexpected upload error".
            msg = f"could not extract text from {fmt_name}: {exc}"
            raise ValueError(msg) from exc

    # Path.suffix is empty for dotfiles (".env") and extensionless names
    # ("Makefile"). Surface the filename in those cases so the failure
    # reason is actionable for the user rather than a naked trailing colon.
    descriptor = ext if ext else f"{path.name} (no extension)"
    msg = f"Unsupported file type: {descriptor}"
    raise ValueError(msg)


def _extract_plain_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


# Per-format extractors stream-and-bail using this cumulative-bytes guard
# so peak memory is bounded too — not just the on-disk store. Without the
# in-loop bail, a heavily-compressed file under the 50 MiB upload cap can
# decompress to hundreds of MB of plain text before _truncate_to_cap fires
# at the boundary. The boundary truncator still applies on top to enforce
# the exact byte cap.
def _extract_pdf(path: Path) -> str:
    pages: list[str] = []
    cumulative = 0
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
                cumulative += len(text.encode("utf-8")) + 2  # "\n\n" separator
                if cumulative >= MAX_EXTRACTED_TEXT_BYTES:
                    break
    return "\n\n".join(pages)


def _extract_docx(path: Path) -> str:
    doc = DocxDocument(str(path))
    parts: list[str] = []
    cumulative = 0
    for p in doc.paragraphs:
        if not p.text:
            continue
        parts.append(p.text)
        cumulative += len(p.text.encode("utf-8")) + 2  # "\n\n" separator
        if cumulative >= MAX_EXTRACTED_TEXT_BYTES:
            break
    return "\n\n".join(parts)


def _extract_pptx(path: Path) -> str:
    prs = PptxPresentation(str(path))
    parts: list[str] = []
    cumulative = 0
    for i, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        # Stream shapes inside the slide and bail mid-slide once the cap
        # is reached (S3). Without this, a slide with thousands of text
        # frames builds up a large in-memory list before the per-slide
        # check runs — same class as the xlsx row-materialisation bug.
        slide_overflow = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                slide_texts.append(text)
                cumulative += len(text.encode("utf-8")) + 1  # "\n" separator
                if cumulative >= MAX_EXTRACTED_TEXT_BYTES:
                    slide_overflow = True
                    break
        if slide_texts:
            joined = f"--- Slide {i} ---\n" + "\n".join(slide_texts)
            parts.append(joined)
            cumulative += len(f"--- Slide {i} ---\n".encode()) + 2  # header + "\n\n"
        if slide_overflow or cumulative >= MAX_EXTRACTED_TEXT_BYTES:
            break
    return "\n\n".join(parts)


def _extract_xlsx(path: Path) -> str:
    wb = load_workbook(path, read_only=True, data_only=True)
    try:
        parts: list[str] = []
        cumulative = 0
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows: list[str] = []
            # Stream cells one at a time inside the row, accumulating into
            # ``cells`` and checking ``cumulative`` mid-row (I6). The
            # previous list-comprehension realised every cell of the row
            # before any cap check fired — a single-row workbook with
            # millions of populated cells could allocate hundreds of MB
            # of Python strings before the post-row check ran.
            for row in sheet.iter_rows():
                cells: list[str] = []
                for cell in row:
                    val = str(cell.value) if cell.value is not None else ""
                    cells.append(val)
                    cumulative += len(val.encode("utf-8")) + 1  # "\t" separator
                    if cumulative >= MAX_EXTRACTED_TEXT_BYTES:
                        break
                rows.append("\t".join(cells))
                if cumulative >= MAX_EXTRACTED_TEXT_BYTES:
                    break
            parts.append(f"--- {sheet_name} ---\n" + "\n".join(rows))
            if cumulative >= MAX_EXTRACTED_TEXT_BYTES:
                break
        return "\n\n".join(parts)
    finally:
        wb.close()


def _extract_rtf(path: Path) -> str:
    # RTF files in the wild commonly contain raw 8-bit bytes (CP1252 /
    # Windows-1252), not just ASCII + ``\'XX`` escapes. Reading as utf-8
    # with errors="replace" would substitute U+FFFD for every non-ASCII
    # byte before striprtf could decode it — silent data loss on
    # legitimate non-ASCII content (S1). CP1252 is a single-byte encoding
    # that round-trips any byte to a Unicode code point.
    raw = path.read_text(encoding="cp1252", errors="replace")
    return str(rtf_to_text(raw))


def get_page_count(path: Path) -> int | None:
    """Return page/sheet/slide count, or None for formats without pages."""
    ext = path.suffix.lower()
    if ext == ".pdf":
        with pdfplumber.open(path) as pdf:
            return len(pdf.pages)
    if ext == ".pptx":
        prs = PptxPresentation(str(path))
        return len(prs.slides)
    if ext == ".xlsx":
        wb = load_workbook(path, read_only=True)
        try:
            return len(wb.sheetnames)
        finally:
            wb.close()
    return None
