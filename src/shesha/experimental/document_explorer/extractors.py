"""Text extraction from uploaded documents.

Dispatches to format-specific extractors based on file extension.
"""

from __future__ import annotations

from pathlib import Path

import pdfplumber
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation as PptxPresentation
from striprtf.striprtf import rtf_to_text

# Extensions treated as plain text (read with open())
_PLAIN_TEXT_EXTENSIONS = frozenset(
    {
        ".txt",
        ".md",
        ".csv",
        ".log",
        ".json",
        ".yaml",
        ".yml",
        ".xml",
        ".html",
        ".htm",
        ".ini",
        ".cfg",
        ".toml",
        ".env",
        ".py",
        ".js",
        ".ts",
        ".java",
        ".c",
        ".cpp",
        ".h",
        ".rs",
        ".go",
        ".rb",
        ".sh",
        ".bat",
        ".sql",
        ".r",
        ".tex",
    }
)


def extract_text(path: Path, content_type: str | None = None) -> str:
    """Extract text content from a file.

    *content_type* is accepted for forward-compatibility but dispatch
    is by file extension.  Raises ``ValueError`` for unsupported types.
    """
    ext = path.suffix.lower()

    if ext in _PLAIN_TEXT_EXTENSIONS:
        return _extract_plain_text(path)
    if ext == ".pdf":
        return _extract_pdf(path)
    if ext == ".docx":
        return _extract_docx(path)
    if ext == ".pptx":
        return _extract_pptx(path)
    if ext == ".xlsx":
        return _extract_xlsx(path)
    if ext == ".rtf":
        return _extract_rtf(path)

    msg = f"Unsupported file type: {ext}"
    raise ValueError(msg)


def _extract_plain_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _extract_pdf(path: Path) -> str:
    pages: list[str] = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
    return "\n\n".join(pages)


def _extract_docx(path: Path) -> str:
    doc = DocxDocument(str(path))
    return "\n\n".join(p.text for p in doc.paragraphs if p.text)


def _extract_pptx(path: Path) -> str:
    prs = PptxPresentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_texts.append(shape.text_frame.text)
        if slide_texts:
            parts.append(f"--- Slide {i} ---\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)


def _extract_xlsx(path: Path) -> str:
    wb = load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        rows: list[str] = []
        for row in sheet.iter_rows():
            cells = [str(cell.value) if cell.value is not None else "" for cell in row]
            rows.append("\t".join(cells))
        parts.append(f"--- {sheet_name} ---\n" + "\n".join(rows))
    return "\n\n".join(parts)


def _extract_rtf(path: Path) -> str:
    raw = path.read_text(encoding="utf-8", errors="replace")
    return str(rtf_to_text(raw))


def get_page_count(path: Path) -> int | None:
    """Return page/sheet/slide count, or None for formats without pages."""
    ext = path.suffix.lower()
    if ext == ".pdf":
        with pdfplumber.open(path) as pdf:
            return len(pdf.pages)
    if ext == ".pptx":
        prs = PptxPresentation(str(path))
        return len(prs.slides)
    if ext == ".xlsx":
        wb = load_workbook(path, read_only=True)
        return len(wb.sheetnames)
    return None
